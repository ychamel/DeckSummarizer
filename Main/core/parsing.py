import uuid
from io import BytesIO
from time import sleep
from typing import List, Any, Optional
import re

import pandas as pd
import streamlit as st
from pypdf import PdfReader
from pptx import Presentation

import docx2txt
from langchain.docstore.document import Document
import fitz
from hashlib import md5

from abc import abstractmethod, ABC
from copy import deepcopy

from Main.core.PDF_Parser import parse_img, fetch_text
from Main.core.WebScrapper import run_spider


class File(ABC):
    """Represents an uploaded file comprised of Documents"""

    def __init__(
            self,
            name: str,
            id: str,
            metadata: Optional[dict[str, Any]] = None,
            docs: Optional[List[Document]] = None,
    ):
        self.name = name
        self.id = id
        self.metadata = metadata or {}
        self.docs = docs or []

    @classmethod
    def from_bytes(cls, file: BytesIO) -> "File":
        """Creates a File from a BytesIO object"""
        return None
    @classmethod
    def from_url(cls, url: str) -> "File":
        """Creates a File from a BytesIO object"""
        return None
    def __repr__(self) -> str:
        return (
            f"File(name={self.name}, id={self.id},"
            " metadata={self.metadata}, docs={self.docs})"
        )

    def __str__(self) -> str:
        return f"File(name={self.name}, id={self.id}, metadata={self.metadata})"

    def copy(self) -> "File":
        """Create a deep copy of this File"""
        return self.__class__(
            name=self.name,
            id=self.id,
            metadata=deepcopy(self.metadata),
            docs=deepcopy(self.docs),
        )


def strip_consecutive_newlines(text: str) -> str:
    """Strips consecutive newlines from a string
    possibly with whitespace in between
    """
    return re.sub(r"\s*\n\s*", "\n", text)


class DocxFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "DocxFile":
        text = docx2txt.process(file)
        text = strip_consecutive_newlines(text)
        doc = Document(page_content=text.strip())
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=[doc])


class PdfFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "PdfFile":
        reader = PdfReader(file)
        docs = []
        uuids = {}
        parsing_bar = st.progress(0.0, text="progress")

        size = len(reader.pages)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            text = strip_consecutive_newlines(text)
            # check ocr enabled
            if st.session_state["OCR_ENABLED"]:
                for image_file_object in page.images:
                    name = image_file_object.name
                    response = parse_img(name, image_file_object.data)
                    if response:
                        uuids[response['uid']] = len(docs)
                    sleep(5)
            doc = Document(page_content=text.strip())
            doc.metadata["page"] = i + 1
            docs.append(doc)
            # update progress
            parsing_bar.progress(i / size, "Parsing PDF")
        parsing_bar.progress(1.0, "Parsing PDF")
        # retrieve images
        progress_text = "Decoding Images"
        # fetching responses
        count = 0
        while len(uuids.keys()) > 0 and count < 10:
            progress = 0.0
            size = len(uuids.keys())

            for uuid, id in list(uuids.items()):
                response = fetch_text(uuid)
                if response['completed']:
                    docs[id].page_content += f" ----- img_data ----- \n {response['document_text']} \n ----- end ----- "
                    del uuids[uuid]
                sleep(1)

                # progress
                progress += 1.0
                parsing_bar.progress(id / size, "Decoding Images")

            # add timeout
            count += 1
            sleep(10)
        parsing_bar.progress(1.0, "Decoding Images")
        # file.read() mutates the file object, which can affect caching
        # so we need to reset the file pointer to the beginning
        file.seek(0)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


class PdfFile2(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "PdfFile2":
        pdf = fitz.open(stream=file.read(), filetype="pdf")  # type: ignore
        docs = []
        uuids = {}
        parsing_bar = st.progress(0.0, text="progress")
        size = len(pdf)
        for i, page in enumerate(pdf):
            text = page.get_text(sort=True)
            text = strip_consecutive_newlines(text)
            # check ocr enabled
            if st.session_state["OCR_ENABLED"]:
                for images in page.get_images():
                    xref = images[0]
                    img = pdf.extract_image(xref)
                    name = str(xref)
                    binary = img['image']
                    response = parse_img(name, binary)
                    if response:
                        uuids[response['uid']] = len(docs)
                    sleep(5)
            doc = Document(page_content=text.strip())
            doc.metadata["page"] = i + 1
            docs.append(doc)
            # update progress
            parsing_bar.progress(i / size, "Parsing PDF")
        # retrieve images
        progress_text = "Decoding Images"
        # fetching responses
        count = 0
        while len(uuids.keys()) > 0 and count < 10:
            progress = 0.0
            size = len(uuids.keys())

            for uuid, id in list(uuids.items()):
                response = fetch_text(uuid)
                if response['completed']:
                    docs[id].page_content += f" ----- img_data ----- \n {response['document_text']} \n ----- end ----- "
                    del uuids[uuid]
                sleep(1)

                # progress
                progress += 1.0
                parsing_bar.progress(id / size, "Decoding Images")

            # add timeout
            count += 1
            sleep(5)
        # file.read() mutates the file object, which can affect caching
        # so we need to reset the file pointer to the beginning
        file.seek(0)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


class TxtFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "TxtFile":
        text = file.read().decode("utf-8")
        text = strip_consecutive_newlines(text)
        file.seek(0)
        doc = Document(page_content=text.strip())
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=[doc])


class XLFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "XLFile":
        dataframes = pd.read_excel(file, None)
        docs = []
        for title, dataframe in dataframes.items():
            dataframe = strip_consecutive_newlines(dataframe.to_string())
            file.seek(0)
            doc = Document(page_content=dataframe.strip())
            docs.append(doc)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


class PPTFile(File):
    @classmethod
    def from_bytes(cls, file: BytesIO) -> "PPTFile":
        # read file
        prs = Presentation(file)

        docs = []
        # loop on slides
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        doc = Document(page_content=run.text.strip())
                        doc.metadata["page"] = i + 1
                        docs.append(doc)
        file.seek(0)
        return cls(name=file.name, id=md5(file.read()).hexdigest(), docs=docs)


def read_file(file: BytesIO) -> File:
    """Reads an uploaded file and returns a File object"""
    if file.name.lower().endswith(".docx"):
        return DocxFile.from_bytes(file)
    elif file.name.lower().endswith(".pdf"):
        return PdfFile.from_bytes(file)
    elif file.name.lower().endswith(".txt"):
        return TxtFile.from_bytes(file)
    elif file.name.lower().endswith(".pptx"):
        return PPTFile.from_bytes(file)
    elif file.name.lower().endswith(".xlsx"):
        return XLFile.from_bytes(file)
    else:
        raise NotImplementedError(f"File type {file.name.split('.')[-1]} not supported")


def scrape_url(url: str) -> List[File]:
    # scrape url
    DICT = run_spider(url)
    # return into files
    files = []
    for web_url, web_content in DICT.items():
        doc = Document(page_content=web_content)
        file = File(name=web_url, id=str(uuid.uuid4()), docs=[doc])
        files.append(file)
    return files
